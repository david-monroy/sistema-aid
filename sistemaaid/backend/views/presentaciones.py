from distutils import text_file
from lib2to3.pgen2.pgen import DFAState
from pptx import Presentation
from django.http import HttpResponse
import pandas as pd
from django.views.decorators.csrf import csrf_exempt
from backend.models import Pregunta,Respuesta, Estudio, Edicion
import json
from django.core.serializers.json import DjangoJSONEncoder
from django.db.models import F
import os
from pptx.chart.data import CategoryChartData,ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt
from django.db.models import Count
from datetime import datetime


@csrf_exempt
def crearPresentacion(request):
    try:
        data = json.loads(request.body.decode("utf-8"))
        prs = Presentation('backend/ia/presentaciones/prueba1.pptx')
        #prs = Presentation()

        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)

        estudioData = data['estudio']
        edicionesData = data['ediciones']
        primero = 1
        ediciones = ''

        # Guarda codigos de ediciones en un array para luego mostrarlos como subtitulo
        for edicion in edicionesData:
            if (primero == 1):
                ediciones = edicion['codigo']
                primero = 0
            else:
                ediciones = ediciones + ', ' + edicion['codigo']
               
        slide.shapes.title.text = 'Reporte de estudio ' + estudioData['nombre'] + ' - ' + estudioData['codigo']
        slide.placeholders[1].text = 'Ediciones: ' + ediciones

        # Agregar logo
        logo = 'backend/assets/logo_AID.png'
        from_left = Inches(5)
        from_top = Inches(2)
        width = Inches(3)
        add_picture = slide.shapes.add_picture(logo, from_left, from_top, width)

        # Agregar fecha
        text_left, text_top, text_width, text_height = Inches(5.7), Inches(6), Inches(2), Inches(1)
        fechaItem = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
        fechaTexto = fechaItem.text_frame.add_paragraph()
        fechaTexto.text = "Fecha: " + datetime.today().strftime('%d-%m-%Y')
        fechaTexto.font.size = Pt(14)

        i = 1
        for item in data['presentacion']: 
            print(item)
            slide_layout = prs.slide_layouts[5]
            print("entro1")
            slide = prs.slides.add_slide(slide_layout)
            print("entro2")
            preguntaData = item["pregunta"]
            print("entro3")
            pregunta = Pregunta.objects.get(pk=preguntaData["id"])
            print("entro4")
            slide.shapes.title.text = pregunta.etiqueta
            print("entro5")
            respuestasData = []
            numero_ediciones = len(item["ediciones"])
            for edicion in item["ediciones"]:
                respuestas = Respuesta.objects.filter(pregunta__edicion__pk=edicion, pregunta__pregunta__pk=preguntaData["id"]).values(
                "respuesta").annotate(Count("id"))
                respuestasData.append(respuestas)
                df = pd.DataFrame(respuestas)

                posicion = item["ediciones"].index(edicion)+1

                if (posicion == 1):
                    x = Inches(7/posicion)
                else:
                    x = Inches((7/posicion)-(0.5*numero_ediciones))

                y, cx, cy = Inches(2.5), Inches(8/numero_ediciones), Inches(7/numero_ediciones)

                codigo_edicion = Edicion.objects.get(pk=edicion)
                etiqueta_series = 'Edición: ' + str(codigo_edicion)

                if (item["tipoGrafico"] == "bar"):
                    print("entro")
                    chart_data = CategoryChartData()
                    chart_data.categories = df["respuesta"]
                    chart_data.add_series(etiqueta_series, df['id__count'])
                    slide.shapes.add_chart(
                    XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
                    )

                if (item["tipoGrafico"] == "pie"):
                    chart_data = ChartData()
                    chart_data.categories = (df["respuesta"])
                    chart_data.add_series(etiqueta_series, df['id__count'])
                    slide.shapes.add_chart(
                    XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
                    )
            
            
            
            print("entro7")
            
            
            i+=1

        prs.save('backend/ia/presentaciones/pruebapresentaciones.pptx')

        return HttpResponse('creó la presentación')

    except BaseException as err:
        return HttpResponse(err)