from pptx import Presentation

prs = Presentation('backend/ia/presentaciones/prueba1.pptx')
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = 'New Title'
prs.save('backend/ia/presentaciones/pruebapresentaciones.pptx')